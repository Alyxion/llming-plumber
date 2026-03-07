"""Build complex PowerPoint presentations from structured JSON definitions.

Supports multiple slide layouts, text boxes, shapes, tables,
images, charts placeholders, speaker notes, and transitions.
"""

from __future__ import annotations

import base64
import io
from typing import Any, ClassVar, Literal

from pydantic import BaseModel, Field

from llming_plumber.blocks.base import BaseBlock, BlockContext, BlockInput, BlockOutput
from llming_plumber.blocks.limits import MAX_SLIDES, check_page_count


class TextStyle(BaseModel):
    """Text formatting options."""

    bold: bool = False
    italic: bool = False
    underline: bool = False
    font_name: str = Field(default="", description="Font name, e.g. 'Calibri'")
    font_size: int | None = Field(
        default=None, description="Font size in points"
    )
    font_color: str = Field(
        default="",
        description="Hex color without #, e.g. 'FF0000'",
    )
    alignment: Literal["left", "center", "right", ""] = ""


class TextRun(BaseModel):
    """A run of text with formatting."""

    text: str
    style: TextStyle = Field(default_factory=TextStyle)
    is_new_paragraph: bool = Field(
        default=False,
        description="Start a new paragraph before this run",
    )


class TextBoxDef(BaseModel):
    """A positioned text box on the slide."""

    left: float = Field(description="Left position in inches")
    top: float = Field(description="Top position in inches")
    width: float = Field(description="Width in inches")
    height: float = Field(description="Height in inches")
    runs: list[TextRun] = Field(description="Formatted text runs")
    bg_color: str = Field(
        default="",
        description="Background fill hex color. Empty for transparent.",
    )


class ShapeDef(BaseModel):
    """A shape element on the slide."""

    shape_type: Literal[
        "rectangle", "rounded_rectangle", "oval",
        "triangle", "arrow_right", "arrow_left",
    ] = "rectangle"
    left: float = Field(description="Left position in inches")
    top: float = Field(description="Top position in inches")
    width: float = Field(description="Width in inches")
    height: float = Field(description="Height in inches")
    fill_color: str = Field(default="", description="Fill hex color")
    line_color: str = Field(default="000000", description="Border hex color")
    text: str = Field(default="", description="Text inside the shape")
    text_style: TextStyle = Field(default_factory=TextStyle)


class SlideImageDef(BaseModel):
    """An image on a slide."""

    data_b64: str = Field(description="Base64-encoded image (PNG or JPEG)")
    left: float = Field(description="Left position in inches")
    top: float = Field(description="Top position in inches")
    width: float = Field(description="Width in inches")
    height: float = Field(description="Height in inches")


class SlideTableDef(BaseModel):
    """A table on a slide."""

    left: float = Field(description="Left position in inches")
    top: float = Field(description="Top position in inches")
    width: float = Field(description="Table width in inches")
    height: float = Field(description="Table height in inches")
    headers: list[str] = Field(default=[], description="Column headers")
    rows: list[list[str]] = Field(description="Data rows")
    header_bg: str = Field(
        default="", description="Header row background hex color"
    )


class SlideDef(BaseModel):
    """Definition of a single slide."""

    layout: Literal[
        "title", "title_and_content", "section_header",
        "two_content", "blank", "title_only",
    ] = Field(
        default="blank",
        description="Slide layout type",
    )
    title: str = Field(default="", description="Slide title text")
    subtitle: str = Field(
        default="",
        description="Subtitle (title slide) or body text (content slides)",
    )
    bullet_points: list[str] = Field(
        default=[],
        description="Bullet points for content placeholder",
    )
    text_boxes: list[TextBoxDef] = Field(
        default=[],
        description="Additional positioned text boxes",
    )
    shapes: list[ShapeDef] = Field(
        default=[],
        description="Shape elements on the slide",
    )
    images: list[SlideImageDef] = Field(
        default=[],
        description="Images on the slide",
    )
    tables: list[SlideTableDef] = Field(
        default=[],
        description="Tables on the slide",
    )
    notes: str = Field(
        default="",
        description="Speaker notes for this slide",
    )
    bg_color: str = Field(
        default="",
        description="Slide background hex color. Empty for default.",
    )


class PowerpointBuilderInput(BlockInput):
    slides: list[SlideDef] = Field(
        title="Slides",
        description="Ordered list of slide definitions",
        min_length=1,
    )
    slide_width: float = Field(
        default=13.333,
        title="Slide Width",
        description="Presentation width in inches (13.333 = widescreen 16:9)",
    )
    slide_height: float = Field(
        default=7.5,
        title="Slide Height",
        description="Presentation height in inches (7.5 = widescreen 16:9)",
    )
    json_data: str = Field(
        default="",
        title="JSON Data",
        description=(
            "Alternative: provide slides as JSON matching SlideDef. "
            "Ignored if slides has data."
        ),
        json_schema_extra={"widget": "code", "rows": 20},
    )


class PowerpointBuilderOutput(BlockOutput):
    content: str = Field(description="Base64-encoded pptx file")
    slide_count: int
    element_count: int


_LAYOUT_MAP: dict[str, int] = {
    "title": 0,
    "title_and_content": 1,
    "section_header": 2,
    "two_content": 3,
    "blank": 6,
    "title_only": 5,
}


class PowerpointBuilderBlock(
    BaseBlock[PowerpointBuilderInput, PowerpointBuilderOutput]
):
    block_type: ClassVar[str] = "powerpoint_builder"
    icon: ClassVar[str] = "tabler/presentation"
    categories: ClassVar[list[str]] = ["documents", "powerpoint"]
    description: ClassVar[str] = (
        "Build complex PowerPoint presentations with text boxes, "
        "shapes, tables, images, and speaker notes"
    )
    cache_ttl: ClassVar[int] = 0

    async def execute(
        self,
        input: PowerpointBuilderInput,
        ctx: BlockContext | None = None,
    ) -> PowerpointBuilderOutput:
        import json

        slides = input.slides
        if input.json_data and not any(
            s.title or s.text_boxes for s in slides
        ):
            raw = json.loads(input.json_data)
            slides = [SlideDef.model_validate(s) for s in raw]

        check_page_count(len(slides), limit=MAX_SLIDES, label="PowerPoint slides")
        return self._build_presentation(
            slides, input.slide_width, input.slide_height
        )

    def _build_presentation(
        self,
        slides: list[SlideDef],
        slide_width: float,
        slide_height: float,
    ) -> PowerpointBuilderOutput:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(slide_width)
        prs.slide_height = Inches(slide_height)

        total_elements = 0

        for slide_def in slides:
            layout_idx = _LAYOUT_MAP.get(slide_def.layout, 6)
            layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(layout)

            if slide_def.bg_color:
                self._set_slide_bg(slide, slide_def.bg_color)

            if slide_def.title and slide.shapes.title:
                slide.shapes.title.text = slide_def.title
                total_elements += 1

            if slide_def.subtitle and slide_def.layout == "title":
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = slide_def.subtitle
                    total_elements += 1

            if slide_def.bullet_points and slide_def.layout in (
                "title_and_content", "two_content"
            ):
                if len(slide.placeholders) > 1:
                    tf = slide.placeholders[1].text_frame
                    for i, point in enumerate(slide_def.bullet_points):
                        if i == 0:
                            tf.text = point
                        else:
                            tf.add_paragraph().text = point
                    total_elements += len(slide_def.bullet_points)

            for tb in slide_def.text_boxes:
                self._add_text_box(slide, tb, Inches, Pt)
                total_elements += 1

            for shape in slide_def.shapes:
                self._add_shape(slide, shape, Inches, Pt)
                total_elements += 1

            for img in slide_def.images:
                self._add_image(slide, img, Inches)
                total_elements += 1

            for table in slide_def.tables:
                self._add_table(slide, table, Inches)
                total_elements += 1

            if slide_def.notes:
                slide.notes_slide.notes_text_frame.text = slide_def.notes
                total_elements += 1

        buf = io.BytesIO()
        prs.save(buf)
        encoded = base64.b64encode(buf.getvalue()).decode("ascii")

        return PowerpointBuilderOutput(
            content=encoded,
            slide_count=len(slides),
            element_count=total_elements,
        )

    def _set_slide_bg(self, slide: Any, hex_color: str) -> None:
        from pptx.dml.color import RGBColor

        bg = slide.background
        fill = bg.fill
        fill.solid()
        h = hex_color
        fill.fore_color.rgb = RGBColor(
            int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
        )

    def _add_text_box(
        self, slide: Any, tb: TextBoxDef, Inches: Any, Pt: Any
    ) -> None:
        from pptx.dml.color import RGBColor

        txBox = slide.shapes.add_textbox(
            Inches(tb.left), Inches(tb.top),
            Inches(tb.width), Inches(tb.height),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        current_para = tf.paragraphs[0]

        for i, run_def in enumerate(tb.runs):
            if run_def.is_new_paragraph and i > 0:
                current_para = tf.add_paragraph()

            run = current_para.add_run()
            run.text = run_def.text
            self._apply_text_style(run, run_def.style, Pt, RGBColor)

            if run_def.style.alignment:
                from pptx.enum.text import PP_ALIGN
                align_map = {
                    "left": PP_ALIGN.LEFT,
                    "center": PP_ALIGN.CENTER,
                    "right": PP_ALIGN.RIGHT,
                }
                current_para.alignment = align_map.get(
                    run_def.style.alignment
                )

        if tb.bg_color:
            fill = txBox.fill
            fill.solid()
            h = tb.bg_color
            fill.fore_color.rgb = RGBColor(
                int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            )

    def _add_shape(
        self, slide: Any, shape: ShapeDef, Inches: Any, Pt: Any
    ) -> None:
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE

        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "arrow_right": MSO_SHAPE.RIGHT_ARROW,
            "arrow_left": MSO_SHAPE.LEFT_ARROW,
        }
        mso_shape = shape_map.get(
            shape.shape_type, MSO_SHAPE.RECTANGLE
        )

        auto_shape = slide.shapes.add_shape(
            mso_shape,
            Inches(shape.left), Inches(shape.top),
            Inches(shape.width), Inches(shape.height),
        )

        if shape.fill_color:
            fill = auto_shape.fill
            fill.solid()
            h = shape.fill_color
            fill.fore_color.rgb = RGBColor(
                int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            )

        if shape.line_color:
            h = shape.line_color
            auto_shape.line.color.rgb = RGBColor(
                int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            )

        if shape.text:
            tf = auto_shape.text_frame
            tf.text = shape.text
            if tf.paragraphs:
                run = tf.paragraphs[0].runs[0]
                self._apply_text_style(
                    run, shape.text_style, Pt, RGBColor
                )

    def _add_image(
        self, slide: Any, img: SlideImageDef, Inches: Any
    ) -> None:
        img_data = base64.b64decode(img.data_b64)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(img.left), Inches(img.top),
            Inches(img.width), Inches(img.height),
        )

    def _add_table(
        self, slide: Any, table: SlideTableDef, Inches: Any
    ) -> None:
        from pptx.dml.color import RGBColor

        total_rows = len(table.rows) + (1 if table.headers else 0)
        num_cols = len(table.headers) if table.headers else (
            len(table.rows[0]) if table.rows else 0
        )
        if num_cols == 0:
            return

        tbl_shape = slide.shapes.add_table(
            total_rows, num_cols,
            Inches(table.left), Inches(table.top),
            Inches(table.width), Inches(table.height),
        )
        tbl = tbl_shape.table

        row_offset = 0
        if table.headers:
            for col_idx, header in enumerate(table.headers):
                cell = tbl.cell(0, col_idx)
                cell.text = header
                if table.header_bg:
                    h = table.header_bg
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(
                        int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
                    )
            row_offset = 1

        for row_idx, row_data in enumerate(table.rows):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < num_cols:
                    tbl.cell(
                        row_idx + row_offset, col_idx
                    ).text = str(cell_text)

    @staticmethod
    def _apply_text_style(
        run: Any, style: TextStyle, Pt: Any, RGBColor: Any
    ) -> None:
        run.font.bold = style.bold
        run.font.italic = style.italic
        run.font.underline = style.underline
        if style.font_name:
            run.font.name = style.font_name
        if style.font_size:
            run.font.size = Pt(style.font_size)
        if style.font_color:
            h = style.font_color
            run.font.color.rgb = RGBColor(
                int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            )
