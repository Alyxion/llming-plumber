"""Read PowerPoint (pptx) presentations."""

from __future__ import annotations

import base64
import io
from typing import Any, ClassVar

from pydantic import Field

from llming_plumber.blocks.base import BaseBlock, BlockContext, BlockInput, BlockOutput


class PowerpointReaderInput(BlockInput):
    content: str = Field(
        title="Content",
        description="Base64-encoded pptx file bytes",
        json_schema_extra={"widget": "textarea"},
    )


class PowerpointReaderOutput(BlockOutput):
    text: str
    slides: list[dict[str, Any]]
    slide_count: int
    metadata: dict[str, Any]


class PowerpointReaderBlock(BaseBlock[PowerpointReaderInput, PowerpointReaderOutput]):
    block_type: ClassVar[str] = "powerpoint_reader"
    icon: ClassVar[str] = "tabler/presentation"
    categories: ClassVar[list[str]] = ["documents", "powerpoint"]
    description: ClassVar[str] = "Read PowerPoint (pptx) presentations"
    cache_ttl: ClassVar[int] = 0

    async def execute(
        self, input: PowerpointReaderInput, ctx: BlockContext | None = None
    ) -> PowerpointReaderOutput:
        from pptx import Presentation

        raw = base64.b64decode(input.content)
        prs = Presentation(io.BytesIO(raw))

        all_text_parts: list[str] = []
        slides_data: list[dict[str, Any]] = []

        for idx, slide in enumerate(prs.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

            notes = ""
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()

            all_text_parts.extend(texts)
            slides_data.append({
                "slide_number": idx,
                "texts": texts,
                "notes": notes,
            })

        metadata: dict[str, Any] = {}
        props = prs.core_properties
        for attr in ("author", "title", "subject", "keywords", "category"):
            val = getattr(props, attr, None)
            if val:
                metadata[attr] = str(val)
        if props.created:
            metadata["created"] = props.created.isoformat()
        if props.modified:
            metadata["modified"] = props.modified.isoformat()

        return PowerpointReaderOutput(
            text="\n".join(all_text_parts),
            slides=slides_data,
            slide_count=len(slides_data),
            metadata=metadata,
        )
