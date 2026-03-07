"""Write PowerPoint (pptx) presentations."""

from __future__ import annotations

import base64
import io
from typing import ClassVar

from pydantic import Field

from llming_plumber.blocks.base import BaseBlock, BlockContext, BlockInput, BlockOutput


class PowerpointWriterInput(BlockInput):
    slides: list[dict[str, str]] = Field(
        title="Slides",
        description="List of slide dicts with 'title' and 'content' strings",
    )
    title: str = Field(
        default="",
        title="Title",
        description="Optional presentation title (added as a title slide)",
    )


class PowerpointWriterOutput(BlockOutput):
    content: str
    slide_count: int


class PowerpointWriterBlock(BaseBlock[PowerpointWriterInput, PowerpointWriterOutput]):
    block_type: ClassVar[str] = "powerpoint_writer"
    icon: ClassVar[str] = "tabler/presentation"
    categories: ClassVar[list[str]] = ["documents", "powerpoint"]
    description: ClassVar[str] = "Write PowerPoint (pptx) presentations"
    cache_ttl: ClassVar[int] = 0

    async def execute(
        self, input: PowerpointWriterInput, ctx: BlockContext | None = None
    ) -> PowerpointWriterOutput:
        from pptx import Presentation

        prs = Presentation()
        slide_count = 0

        if input.title:
            title_layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = input.title
            slide_count += 1

        content_layout = prs.slide_layouts[1]  # Title and Content
        for slide_data in input.slides:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide.placeholders[1]
            body.text = slide_data.get("content", "")
            slide_count += 1

        buf = io.BytesIO()
        prs.save(buf)
        encoded = base64.b64encode(buf.getvalue()).decode("ascii")

        return PowerpointWriterOutput(
            content=encoded,
            slide_count=slide_count,
        )
