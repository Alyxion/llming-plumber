"""Extract structured SlideDef models from PowerPoint files.

Outputs the same Pydantic models that PowerpointBuilderBlock consumes,
enabling round-trip: extract → edit JSON → rebuild.
"""

from __future__ import annotations

import base64
import io
from typing import Any, ClassVar

from pydantic import Field

from llming_plumber.blocks.base import (
    BaseBlock,
    BlockContext,
    BlockInput,
    BlockOutput,
)
from llming_plumber.blocks.documents.powerpoint_builder import (
    ShapeDef,
    SlideDef,
    SlideTableDef,
    TextBoxDef,
    TextRun,
    TextStyle,
)


class PowerpointExtractorInput(BlockInput):
    content: str = Field(
        title="Content",
        description="Base64-encoded pptx file bytes",
        json_schema_extra={"widget": "textarea"},
    )
    extract_shapes: bool = Field(
        default=True,
        title="Extract Shapes",
        description="Include shape elements (rectangles, etc.)",
    )
    extract_images: bool = Field(
        default=False,
        title="Extract Images",
        description=(
            "Include images as base64. "
            "Warning: increases output size significantly."
        ),
    )


class PowerpointExtractorOutput(BlockOutput):
    slides: list[SlideDef]
    slides_json: str = Field(
        description=(
            "JSON string of slides for piping to builders"
        ),
    )


class PowerpointExtractorBlock(
    BaseBlock[
        PowerpointExtractorInput, PowerpointExtractorOutput
    ]
):
    block_type: ClassVar[str] = "powerpoint_extractor"
    icon: ClassVar[str] = "tabler/presentation"
    categories: ClassVar[list[str]] = [
        "documents", "powerpoint",
    ]
    description: ClassVar[str] = (
        "Extract structured SlideDef models from PowerPoint "
        "for editing and rebuilding"
    )
    cache_ttl: ClassVar[int] = 0

    async def execute(
        self,
        input: PowerpointExtractorInput,
        ctx: BlockContext | None = None,
    ) -> PowerpointExtractorOutput:
        import json

        from pptx import Presentation
        from pptx.util import Emu

        raw = base64.b64decode(input.content)
        prs = Presentation(io.BytesIO(raw))

        slides: list[SlideDef] = []

        for slide in prs.slides:
            slides.append(
                self._extract_slide(
                    slide,
                    input.extract_shapes,
                    input.extract_images,
                    Emu,
                )
            )

        slides_json = json.dumps(
            [s.model_dump() for s in slides], default=str
        )

        return PowerpointExtractorOutput(
            slides=slides, slides_json=slides_json
        )

    def _extract_slide(
        self,
        slide: Any,
        extract_shapes: bool,
        extract_images: bool,
        Emu: Any,
    ) -> SlideDef:
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text

        notes = ""
        if (
            slide.has_notes_slide
            and slide.notes_slide.notes_text_frame
        ):
            notes = (
                slide.notes_slide.notes_text_frame.text.strip()
            )

        text_boxes: list[TextBoxDef] = []
        shapes: list[ShapeDef] = []
        tables: list[SlideTableDef] = []
        bullet_points: list[str] = []

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue

            if shape.has_table:
                tables.append(
                    self._extract_table(shape, Emu)
                )
            elif shape.has_text_frame:
                if self._is_placeholder(shape):
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            bullet_points.append(para.text)
                else:
                    tb = self._extract_text_box(shape, Emu)
                    if tb:
                        text_boxes.append(tb)
            elif extract_shapes and shape.shape_type is not None:
                s = self._extract_shape(shape, Emu)
                if s:
                    shapes.append(s)

        # Infer layout from content
        if title and bullet_points:
            layout = "title_and_content"
        elif title:
            layout = "title_only"
        else:
            layout = "blank"

        return SlideDef(
            layout=layout,
            title=title,
            bullet_points=bullet_points,
            text_boxes=text_boxes,
            shapes=shapes,
            tables=tables,
            notes=notes,
        )

    def _is_placeholder(self, shape: Any) -> bool:
        try:
            return shape.placeholder_format is not None
        except ValueError:
            return False

    def _emu_to_inches(self, emu_val: Any) -> float:
        if emu_val is None:
            return 0
        return round(int(emu_val) / 914400, 3)

    def _extract_text_box(
        self, shape: Any, Emu: Any
    ) -> TextBoxDef | None:
        runs: list[TextRun] = []
        first_para = True

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text:
                    continue

                style = TextStyle()
                if run.font.bold:
                    style = TextStyle(
                        bold=True,
                        italic=style.italic,
                        font_name=style.font_name,
                        font_size=style.font_size,
                        font_color=style.font_color,
                    )
                if run.font.italic:
                    style = TextStyle(
                        bold=style.bold,
                        italic=True,
                        font_name=style.font_name,
                        font_size=style.font_size,
                        font_color=style.font_color,
                    )
                if run.font.name:
                    style = TextStyle(
                        bold=style.bold,
                        italic=style.italic,
                        font_name=run.font.name,
                        font_size=style.font_size,
                        font_color=style.font_color,
                    )
                if run.font.size:
                    style = TextStyle(
                        bold=style.bold,
                        italic=style.italic,
                        font_name=style.font_name,
                        font_size=int(
                            run.font.size.pt
                            if hasattr(run.font.size, "pt")
                            else run.font.size
                        ),
                        font_color=style.font_color,
                    )
                try:
                    rgb = (
                        run.font.color.rgb
                        if run.font.color
                        else None
                    )
                except AttributeError:
                    rgb = None
                if rgb:
                    style = TextStyle(
                        bold=style.bold,
                        italic=style.italic,
                        font_name=style.font_name,
                        font_size=style.font_size,
                        font_color=str(rgb),
                    )

                runs.append(TextRun(
                    text=run.text,
                    style=style,
                    is_new_paragraph=not first_para,
                ))

            first_para = False

        if not runs:
            return None

        return TextBoxDef(
            left=self._emu_to_inches(shape.left),
            top=self._emu_to_inches(shape.top),
            width=self._emu_to_inches(shape.width),
            height=self._emu_to_inches(shape.height),
            runs=runs,
        )

    def _extract_shape(
        self, shape: Any, Emu: Any
    ) -> ShapeDef | None:
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text

        return ShapeDef(
            shape_type="rectangle",
            left=self._emu_to_inches(shape.left),
            top=self._emu_to_inches(shape.top),
            width=self._emu_to_inches(shape.width),
            height=self._emu_to_inches(shape.height),
            text=text,
        )

    def _extract_table(
        self, shape: Any, Emu: Any
    ) -> SlideTableDef:
        tbl = shape.table
        headers: list[str] = []
        rows: list[list[str]] = []

        for row_idx in range(len(tbl.rows)):
            row_data = [
                tbl.cell(row_idx, col_idx).text
                for col_idx in range(len(tbl.columns))
            ]
            if row_idx == 0:
                headers = row_data
            else:
                rows.append(row_data)

        return SlideTableDef(
            left=self._emu_to_inches(shape.left),
            top=self._emu_to_inches(shape.top),
            width=self._emu_to_inches(shape.width),
            height=self._emu_to_inches(shape.height),
            headers=headers,
            rows=rows,
        )
