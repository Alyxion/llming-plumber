"""Unit tests for all document processing blocks."""

from __future__ import annotations

import base64
import io
from typing import Any
from unittest.mock import MagicMock, patch

import openpyxl
import pyarrow as pa
import pyarrow.parquet as pq
from docx import Document as DocxDocument
from pptx import Presentation

from llming_plumber.blocks.documents.excel_reader import (
    ExcelReaderBlock,
    ExcelReaderInput,
)
from llming_plumber.blocks.documents.excel_writer import (
    ExcelWriterBlock,
    ExcelWriterInput,
)
from llming_plumber.blocks.documents.markdown_renderer import (
    MarkdownRendererBlock,
    MarkdownRendererInput,
)
from llming_plumber.blocks.documents.parquet_reader import (
    ParquetReaderBlock,
    ParquetReaderInput,
)
from llming_plumber.blocks.documents.parquet_writer import (
    ParquetWriterBlock,
    ParquetWriterInput,
)
from llming_plumber.blocks.documents.pdf_reader import (
    PdfReaderBlock,
    PdfReaderInput,
)
from llming_plumber.blocks.documents.pdf_renderer import (
    PdfRendererBlock,
    PdfRendererInput,
)
from llming_plumber.blocks.documents.powerpoint_reader import (
    PowerpointReaderBlock,
    PowerpointReaderInput,
)
from llming_plumber.blocks.documents.powerpoint_writer import (
    PowerpointWriterBlock,
    PowerpointWriterInput,
)
from llming_plumber.blocks.documents.word_reader import (
    WordReaderBlock,
    WordReaderInput,
)
from llming_plumber.blocks.documents.word_writer import (
    WordWriterBlock,
    WordWriterInput,
)
from llming_plumber.blocks.documents.yaml_parser import (
    YamlParserBlock,
    YamlParserInput,
)
from llming_plumber.blocks.documents.yaml_writer import (
    YamlWriterBlock,
    YamlWriterInput,
)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_xlsx(
    rows: list[list[Any]], sheet_name: str = "Sheet1"
) -> str:
    """Create a minimal xlsx file in memory, return base64-encoded string."""
    wb = openpyxl.Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = sheet_name
    for row in rows:
        ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return base64.b64encode(buf.getvalue()).decode("ascii")


def _make_docx(paragraphs: list[str]) -> str:
    """Create a minimal docx file in memory, return base64-encoded string."""
    doc = DocxDocument()
    for text in paragraphs:
        doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return base64.b64encode(buf.getvalue()).decode("ascii")


def _make_pptx(slides: list[tuple[str, str]]) -> str:
    """Create a minimal pptx file, return base64-encoded string."""
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for title, content in slides:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = content
    buf = io.BytesIO()
    prs.save(buf)
    return base64.b64encode(buf.getvalue()).decode("ascii")


def _make_parquet(records: list[dict[str, Any]]) -> str:
    """Create a minimal parquet file, return base64-encoded string."""
    if not records:
        table = pa.table({})
    else:
        cols = list(records[0].keys())
        columnar = {col: [r.get(col) for r in records] for col in cols}
        table = pa.table(columnar)
    buf = io.BytesIO()
    pq.write_table(table, buf)
    return base64.b64encode(buf.getvalue()).decode("ascii")


def _mock_pdfplumber(
    pages_text: list[str],
    tables: list[list[list[list[str]]]] | None = None,
    metadata: dict[str, str] | None = None,
) -> MagicMock:
    """Build a mock that replaces pdfplumber.open()."""
    mock_pdf = MagicMock()
    mock_pages = []
    for i, text in enumerate(pages_text):
        page = MagicMock()
        page.extract_text.return_value = text
        page_tables = tables[i] if tables else []
        page.extract_tables.return_value = page_tables
        mock_pages.append(page)
    mock_pdf.pages = mock_pages
    mock_pdf.metadata = metadata or {}
    return mock_pdf


# ---------------------------------------------------------------------------
# PDF Reader (mocked)
# ---------------------------------------------------------------------------


async def test_pdf_reader_basic() -> None:
    block = PdfReaderBlock()
    content = base64.b64encode(b"fake-pdf").decode("ascii")
    mock_pdf = _mock_pdfplumber(["Hello world", "Page two"])
    mock_mod = MagicMock()
    mock_mod.open.return_value = mock_pdf
    with patch.dict("sys.modules", {"pdfplumber": mock_mod}):
        result = await block.execute(PdfReaderInput(content=content))

    assert result.page_count == 2
    assert len(result.pages) == 2
    assert result.pages[0]["text"] == "Hello world"
    assert result.pages[1]["text"] == "Page two"
    assert "Hello world\nPage two" == result.text


async def test_pdf_reader_specific_pages() -> None:
    block = PdfReaderBlock()
    content = base64.b64encode(b"fake-pdf").decode("ascii")
    mock_pdf = _mock_pdfplumber(["P1", "P2", "P3"])
    mock_mod = MagicMock()
    mock_mod.open.return_value = mock_pdf
    with patch.dict("sys.modules", {"pdfplumber": mock_mod}):
        result = await block.execute(
            PdfReaderInput(content=content, pages=[2])
        )

    assert len(result.pages) == 1
    assert result.pages[0]["page_number"] == 2
    assert result.pages[0]["text"] == "P2"


async def test_pdf_reader_extract_tables() -> None:
    block = PdfReaderBlock()
    content = base64.b64encode(b"fake-pdf").decode("ascii")
    tables_data: list[list[list[list[str]]]] = [
        [[["A", "B"], ["1", "2"]]],
    ]
    mock_pdf = _mock_pdfplumber(["text"], tables=tables_data)
    mock_mod = MagicMock()
    mock_mod.open.return_value = mock_pdf
    with patch.dict("sys.modules", {"pdfplumber": mock_mod}):
        result = await block.execute(
            PdfReaderInput(content=content, extract_tables=True)
        )

    assert "tables" in result.pages[0]
    assert result.pages[0]["tables"] == tables_data[0]


async def test_pdf_reader_metadata() -> None:
    block = PdfReaderBlock()
    content = base64.b64encode(b"fake-pdf").decode("ascii")
    meta = {"Title": "Test Doc", "Author": "Tester"}
    mock_pdf = _mock_pdfplumber(["hello"], metadata=meta)
    mock_mod = MagicMock()
    mock_mod.open.return_value = mock_pdf
    with patch.dict("sys.modules", {"pdfplumber": mock_mod}):
        result = await block.execute(PdfReaderInput(content=content))

    assert result.metadata["title"] == "Test Doc"
    assert result.metadata["author"] == "Tester"


async def test_pdf_reader_out_of_range_page() -> None:
    block = PdfReaderBlock()
    content = base64.b64encode(b"fake-pdf").decode("ascii")
    mock_pdf = _mock_pdfplumber(["only page"])
    mock_mod = MagicMock()
    mock_mod.open.return_value = mock_pdf
    with patch.dict("sys.modules", {"pdfplumber": mock_mod}):
        result = await block.execute(
            PdfReaderInput(content=content, pages=[5])
        )

    assert len(result.pages) == 0
    assert result.page_count == 1


async def test_pdf_reader_block_type() -> None:
    assert PdfReaderBlock.block_type == "pdf_reader"


# ---------------------------------------------------------------------------
# PDF Renderer (mocked)
# ---------------------------------------------------------------------------


def _make_mock_pdfium(
    num_pages: int = 1,
    width: int = 100,
    height: int = 150,
    image_bytes: bytes = b"fake-image-data",
) -> tuple[MagicMock, MagicMock]:
    """Build a mock pypdfium2 module and its PdfDocument return value."""
    mock_pil_image = MagicMock()
    mock_pil_image.width = width
    mock_pil_image.height = height
    mock_pil_image.save.side_effect = lambda buf, format: buf.write(image_bytes)

    mock_bitmap = MagicMock()
    mock_bitmap.to_pil.return_value = mock_pil_image

    mock_page = MagicMock()
    mock_page.render.return_value = mock_bitmap

    mock_pdf = MagicMock()
    mock_pdf.__len__ = MagicMock(return_value=num_pages)
    mock_pdf.__getitem__ = MagicMock(return_value=mock_page)

    mock_mod = MagicMock()
    mock_mod.PdfDocument.return_value = mock_pdf
    return mock_mod, mock_pil_image


async def test_pdf_renderer_basic() -> None:
    block = PdfRendererBlock()
    content = base64.b64encode(b"fake-pdf").decode("ascii")

    mock_mod, _ = _make_mock_pdfium()
    with patch.dict("sys.modules", {"pypdfium2": mock_mod}):
        result = await block.execute(PdfRendererInput(content=content))

    assert result.page_count == 1
    assert len(result.images) == 1
    assert result.images[0]["width"] == 100
    assert result.images[0]["height"] == 150
    decoded = base64.b64decode(result.images[0]["image"])
    assert decoded == b"fake-image-data"


async def test_pdf_renderer_png_format() -> None:
    block = PdfRendererBlock()
    content = base64.b64encode(b"fake-pdf").decode("ascii")

    mock_mod, mock_pil_image = _make_mock_pdfium(
        width=50, height=75, image_bytes=b"png-data"
    )
    with patch.dict("sys.modules", {"pypdfium2": mock_mod}):
        await block.execute(
            PdfRendererInput(content=content, format="png")
        )

    mock_pil_image.save.assert_called_once()
    call_args = mock_pil_image.save.call_args
    assert call_args[1]["format"] == "PNG" or call_args[0][1] == "PNG"


async def test_pdf_renderer_block_type() -> None:
    assert PdfRendererBlock.block_type == "pdf_renderer"


# ---------------------------------------------------------------------------
# Excel Reader
# ---------------------------------------------------------------------------


async def test_excel_reader_xlsx_basic() -> None:
    block = ExcelReaderBlock()
    content = _make_xlsx([["name", "age"], ["Alice", 30], ["Bob", 25]])
    result = await block.execute(ExcelReaderInput(content=content))

    assert result.columns == ["name", "age"]
    assert result.row_count == 2
    assert result.records[0]["name"] == "Alice"
    assert result.records[1]["age"] == 25


async def test_excel_reader_xlsx_no_header() -> None:
    block = ExcelReaderBlock()
    content = _make_xlsx([["Alice", 30], ["Bob", 25]])
    result = await block.execute(
        ExcelReaderInput(content=content, header_row=0)
    )

    assert result.columns == ["col_0", "col_1"]
    assert result.row_count == 2


async def test_excel_reader_xlsx_named_sheet() -> None:
    block = ExcelReaderBlock()
    wb = openpyxl.Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "Data"
    ws.append(["x", "y"])
    ws.append([1, 2])
    buf = io.BytesIO()
    wb.save(buf)
    content = base64.b64encode(buf.getvalue()).decode("ascii")

    result = await block.execute(
        ExcelReaderInput(content=content, sheet_name="Data")
    )
    assert result.columns == ["x", "y"]
    assert result.records[0]["x"] == 1
    assert "Data" in result.sheet_names


async def test_excel_reader_empty_sheet() -> None:
    block = ExcelReaderBlock()
    content = _make_xlsx([])
    result = await block.execute(ExcelReaderInput(content=content))

    assert result.records == []
    assert result.row_count == 0


async def test_excel_reader_block_type() -> None:
    assert ExcelReaderBlock.block_type == "excel_reader"


# ---------------------------------------------------------------------------
# Excel Writer
# ---------------------------------------------------------------------------


async def test_excel_writer_basic() -> None:
    block = ExcelWriterBlock()
    records = [{"name": "Alice", "age": 30}, {"name": "Bob", "age": 25}]
    result = await block.execute(ExcelWriterInput(records=records))

    assert result.row_count == 2
    assert result.column_count == 2

    # Verify output is valid xlsx
    raw = base64.b64decode(result.content)
    wb = openpyxl.load_workbook(io.BytesIO(raw))
    ws = wb.active
    assert ws is not None
    rows = list(ws.iter_rows(values_only=True))
    assert rows[0] == ("name", "age")
    assert rows[1] == ("Alice", 30)


async def test_excel_writer_empty_records() -> None:
    block = ExcelWriterBlock()
    result = await block.execute(ExcelWriterInput(records=[]))

    assert result.row_count == 0
    assert result.column_count == 0
    # Still valid xlsx
    raw = base64.b64decode(result.content)
    wb = openpyxl.load_workbook(io.BytesIO(raw))
    assert wb.active is not None


async def test_excel_writer_custom_columns() -> None:
    block = ExcelWriterBlock()
    records = [{"a": 1, "b": 2, "c": 3}]
    result = await block.execute(
        ExcelWriterInput(records=records, columns=["c", "a"])
    )

    raw = base64.b64decode(result.content)
    wb = openpyxl.load_workbook(io.BytesIO(raw))
    ws = wb.active
    assert ws is not None
    rows = list(ws.iter_rows(values_only=True))
    assert rows[0] == ("c", "a")
    assert rows[1] == (3, 1)
    assert result.column_count == 2


async def test_excel_writer_block_type() -> None:
    assert ExcelWriterBlock.block_type == "excel_writer"


# ---------------------------------------------------------------------------
# Word Reader
# ---------------------------------------------------------------------------


async def test_word_reader_basic() -> None:
    block = WordReaderBlock()
    content = _make_docx(["Hello", "World"])
    result = await block.execute(WordReaderInput(content=content))

    assert "Hello" in result.text
    assert "World" in result.text
    assert len(result.paragraphs) >= 2
    found_texts = [p["text"] for p in result.paragraphs]
    assert "Hello" in found_texts
    assert "World" in found_texts


async def test_word_reader_with_table() -> None:
    block = WordReaderBlock()
    doc = DocxDocument()
    doc.add_paragraph("Before table")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    buf = io.BytesIO()
    doc.save(buf)
    content = base64.b64encode(buf.getvalue()).decode("ascii")

    result = await block.execute(WordReaderInput(content=content))

    assert len(result.tables) == 1
    assert result.tables[0][0] == ["A", "B"]
    assert result.tables[0][1] == ["1", "2"]


async def test_word_reader_block_type() -> None:
    assert WordReaderBlock.block_type == "word_reader"


# ---------------------------------------------------------------------------
# Word Writer
# ---------------------------------------------------------------------------


async def test_word_writer_basic() -> None:
    block = WordWriterBlock()
    paragraphs = [
        {"text": "First paragraph", "style": "Normal"},
        {"text": "Second paragraph"},
    ]
    result = await block.execute(WordWriterInput(paragraphs=paragraphs))

    assert result.paragraph_count == 2
    # Verify output is valid docx
    raw = base64.b64decode(result.content)
    doc = DocxDocument(io.BytesIO(raw))
    texts = [p.text for p in doc.paragraphs]
    assert "First paragraph" in texts
    assert "Second paragraph" in texts


async def test_word_writer_with_title() -> None:
    block = WordWriterBlock()
    paragraphs = [{"text": "Body text"}]
    result = await block.execute(
        WordWriterInput(paragraphs=paragraphs, title="My Document")
    )

    raw = base64.b64decode(result.content)
    doc = DocxDocument(io.BytesIO(raw))
    # Title heading should be present
    texts = [p.text for p in doc.paragraphs]
    assert "My Document" in texts


async def test_word_writer_block_type() -> None:
    assert WordWriterBlock.block_type == "word_writer"


# ---------------------------------------------------------------------------
# PowerPoint Reader
# ---------------------------------------------------------------------------


async def test_powerpoint_reader_basic() -> None:
    block = PowerpointReaderBlock()
    content = _make_pptx([("Slide 1", "Content A"), ("Slide 2", "Content B")])
    result = await block.execute(PowerpointReaderInput(content=content))

    assert result.slide_count == 2
    assert len(result.slides) == 2
    assert "Slide 1" in result.text
    assert "Content A" in result.text
    assert result.slides[0]["slide_number"] == 1
    assert result.slides[1]["slide_number"] == 2


async def test_powerpoint_reader_block_type() -> None:
    assert PowerpointReaderBlock.block_type == "powerpoint_reader"


# ---------------------------------------------------------------------------
# PowerPoint Writer
# ---------------------------------------------------------------------------


async def test_powerpoint_writer_basic() -> None:
    block = PowerpointWriterBlock()
    slides = [
        {"title": "Intro", "content": "Hello"},
        {"title": "Details", "content": "World"},
    ]
    result = await block.execute(PowerpointWriterInput(slides=slides))

    assert result.slide_count == 2
    # Verify output is valid pptx
    raw = base64.b64decode(result.content)
    prs = Presentation(io.BytesIO(raw))
    assert len(prs.slides) == 2


async def test_powerpoint_writer_with_title_slide() -> None:
    block = PowerpointWriterBlock()
    slides = [{"title": "Body", "content": "text"}]
    result = await block.execute(
        PowerpointWriterInput(slides=slides, title="My Presentation")
    )

    assert result.slide_count == 2  # title slide + 1 content slide
    raw = base64.b64decode(result.content)
    prs = Presentation(io.BytesIO(raw))
    assert len(prs.slides) == 2


async def test_powerpoint_writer_block_type() -> None:
    assert PowerpointWriterBlock.block_type == "powerpoint_writer"


# ---------------------------------------------------------------------------
# YAML Parser
# ---------------------------------------------------------------------------


async def test_yaml_parser_single_document() -> None:
    block = YamlParserBlock()
    result = await block.execute(
        YamlParserInput(content="name: Alice\nage: 30\n")
    )

    assert result.data == {"name": "Alice", "age": 30}
    assert result.document_count == 1


async def test_yaml_parser_nested() -> None:
    block = YamlParserBlock()
    yaml_text = "server:\n  host: localhost\n  port: 8080\n"
    result = await block.execute(YamlParserInput(content=yaml_text))

    assert result.data["server"]["host"] == "localhost"
    assert result.data["server"]["port"] == 8080


async def test_yaml_parser_list() -> None:
    block = YamlParserBlock()
    result = await block.execute(
        YamlParserInput(content="- one\n- two\n- three\n")
    )

    assert result.data == ["one", "two", "three"]


async def test_yaml_parser_multi_document() -> None:
    block = YamlParserBlock()
    yaml_text = "a: 1\n---\nb: 2\n"
    result = await block.execute(
        YamlParserInput(content=yaml_text, multi_document=True)
    )

    assert result.document_count == 2
    assert result.data[0] == {"a": 1}
    assert result.data[1] == {"b": 2}


async def test_yaml_parser_block_type() -> None:
    assert YamlParserBlock.block_type == "yaml_parser"


# ---------------------------------------------------------------------------
# YAML Writer
# ---------------------------------------------------------------------------


async def test_yaml_writer_dict() -> None:
    block = YamlWriterBlock()
    result = await block.execute(
        YamlWriterInput(data={"name": "Alice", "age": 30})
    )

    assert "name: Alice" in result.content
    assert "age: 30" in result.content


async def test_yaml_writer_list() -> None:
    block = YamlWriterBlock()
    result = await block.execute(YamlWriterInput(data=["a", "b", "c"]))

    assert "- a" in result.content
    assert "- b" in result.content


async def test_yaml_writer_flow_style() -> None:
    block = YamlWriterBlock()
    result = await block.execute(
        YamlWriterInput(data={"x": 1}, default_flow_style=True)
    )

    assert "{x: 1}" in result.content


async def test_yaml_writer_block_type() -> None:
    assert YamlWriterBlock.block_type == "yaml_writer"


# ---------------------------------------------------------------------------
# Parquet Reader
# ---------------------------------------------------------------------------


async def test_parquet_reader_basic() -> None:
    block = ParquetReaderBlock()
    records = [{"name": "Alice", "age": 30}, {"name": "Bob", "age": 25}]
    content = _make_parquet(records)
    result = await block.execute(ParquetReaderInput(content=content))

    assert result.row_count == 2
    assert "name" in result.columns
    assert "age" in result.columns
    assert result.records[0]["name"] == "Alice"
    assert result.records[1]["age"] == 25


async def test_parquet_reader_select_columns() -> None:
    block = ParquetReaderBlock()
    records = [{"a": 1, "b": 2, "c": 3}]
    content = _make_parquet(records)
    result = await block.execute(
        ParquetReaderInput(content=content, columns=["a", "c"])
    )

    assert result.columns == ["a", "c"]
    assert "b" not in result.records[0]


async def test_parquet_reader_column_schema() -> None:
    block = ParquetReaderBlock()
    records = [{"x": 1, "y": "hello"}]
    content = _make_parquet(records)
    result = await block.execute(ParquetReaderInput(content=content))

    schema_names = [s["name"] for s in result.column_schema]
    assert "x" in schema_names
    assert "y" in schema_names


async def test_parquet_reader_block_type() -> None:
    assert ParquetReaderBlock.block_type == "parquet_reader"


# ---------------------------------------------------------------------------
# Parquet Writer
# ---------------------------------------------------------------------------


async def test_parquet_writer_basic() -> None:
    block = ParquetWriterBlock()
    records = [{"name": "Alice", "score": 95}, {"name": "Bob", "score": 87}]
    result = await block.execute(ParquetWriterInput(records=records))

    assert result.row_count == 2
    assert result.column_count == 2
    # Verify output is valid parquet
    raw = base64.b64decode(result.content)
    table = pq.read_table(io.BytesIO(raw))
    assert table.num_rows == 2


async def test_parquet_writer_empty_records() -> None:
    block = ParquetWriterBlock()
    result = await block.execute(ParquetWriterInput(records=[]))

    assert result.row_count == 0
    assert result.column_count == 0
    # Still valid parquet
    raw = base64.b64decode(result.content)
    table = pq.read_table(io.BytesIO(raw))
    assert table.num_rows == 0


async def test_parquet_writer_block_type() -> None:
    assert ParquetWriterBlock.block_type == "parquet_writer"


# ---------------------------------------------------------------------------
# Markdown Renderer
# ---------------------------------------------------------------------------


async def test_markdown_renderer_heading() -> None:
    block = MarkdownRendererBlock()
    result = await block.execute(
        MarkdownRendererInput(content="# Hello World")
    )

    assert "<h1>" in result.html
    assert "Hello World" in result.html
    assert "Hello World" in result.text


async def test_markdown_renderer_table() -> None:
    block = MarkdownRendererBlock()
    md_table = "| A | B |\n|---|---|\n| 1 | 2 |"
    result = await block.execute(MarkdownRendererInput(content=md_table))

    assert "<table>" in result.html
    assert "<td>" in result.html


async def test_markdown_renderer_code_block() -> None:
    block = MarkdownRendererBlock()
    md_code = "```python\nprint('hello')\n```"
    result = await block.execute(MarkdownRendererInput(content=md_code))

    assert "<code" in result.html
    assert "print('hello')" in result.text


async def test_markdown_renderer_plain_text_extraction() -> None:
    block = MarkdownRendererBlock()
    result = await block.execute(
        MarkdownRendererInput(content="**bold** and *italic*")
    )

    assert "bold" in result.text
    assert "italic" in result.text
    # Text should not contain HTML tags
    assert "<strong>" not in result.text


async def test_markdown_renderer_block_type() -> None:
    assert MarkdownRendererBlock.block_type == "markdown_renderer"
