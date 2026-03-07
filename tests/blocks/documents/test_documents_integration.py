"""Integration tests for document processing blocks using real fixture files.

No mocks -- these tests exercise the actual parsing/rendering libraries
against real files on disk.
"""

from __future__ import annotations

import base64
import io
from pathlib import Path
from typing import Any

import openpyxl
import pyarrow.parquet as pq
import yaml
from docx import Document as DocxDocument
from pptx import Presentation

from llming_plumber.blocks.documents.excel_reader import (
    ExcelReaderBlock,
    ExcelReaderInput,
)
from llming_plumber.blocks.documents.excel_writer import (
    ExcelWriterBlock,
    ExcelWriterInput,
)
from llming_plumber.blocks.documents.markdown_renderer import (
    MarkdownRendererBlock,
    MarkdownRendererInput,
)
from llming_plumber.blocks.documents.parquet_reader import (
    ParquetReaderBlock,
    ParquetReaderInput,
)
from llming_plumber.blocks.documents.parquet_writer import (
    ParquetWriterBlock,
    ParquetWriterInput,
)
from llming_plumber.blocks.documents.pdf_reader import (
    PdfReaderBlock,
    PdfReaderInput,
)
from llming_plumber.blocks.documents.pdf_renderer import (
    PdfRendererBlock,
    PdfRendererInput,
)
from llming_plumber.blocks.documents.powerpoint_reader import (
    PowerpointReaderBlock,
    PowerpointReaderInput,
)
from llming_plumber.blocks.documents.powerpoint_writer import (
    PowerpointWriterBlock,
    PowerpointWriterInput,
)
from llming_plumber.blocks.documents.word_reader import (
    WordReaderBlock,
    WordReaderInput,
)
from llming_plumber.blocks.documents.word_writer import (
    WordWriterBlock,
    WordWriterInput,
)
from llming_plumber.blocks.documents.yaml_parser import (
    YamlParserBlock,
    YamlParserInput,
)
from llming_plumber.blocks.documents.yaml_writer import (
    YamlWriterBlock,
    YamlWriterInput,
)

FIXTURES = Path(__file__).parent.parent.parent / "fixtures"


def _b64(filename: str) -> str:
    """Load a fixture file and return its content as a base64-encoded string."""
    return base64.b64encode((FIXTURES / filename).read_bytes()).decode()


# ---------------------------------------------------------------------------
# PDF Reader
# ---------------------------------------------------------------------------


async def test_pdf_reader_sample_pdf() -> None:
    block = PdfReaderBlock()
    result = await block.execute(PdfReaderInput(content=_b64("sample.pdf")))

    assert result.page_count == 7
    assert len(result.pages) == 7
    assert len(result.text) > 0


async def test_pdf_reader_sample_text_pdf() -> None:
    block = PdfReaderBlock()
    result = await block.execute(PdfReaderInput(content=_b64("sample_text.pdf")))

    assert result.page_count >= 1
    # SCOTUS transcript should contain recognizable words
    assert len(result.text) > 100


async def test_pdf_reader_specific_page() -> None:
    block = PdfReaderBlock()
    result = await block.execute(
        PdfReaderInput(content=_b64("sample.pdf"), pages=[1])
    )

    assert len(result.pages) == 1
    assert result.pages[0]["page_number"] == 1
    # page_count still reports total pages in the document
    assert result.page_count == 7


async def test_pdf_reader_extract_tables() -> None:
    block = PdfReaderBlock()
    result = await block.execute(
        PdfReaderInput(content=_b64("sample.pdf"), extract_tables=True)
    )

    # At least some pages should have a "tables" key
    pages_with_tables = [p for p in result.pages if "tables" in p]
    assert len(pages_with_tables) > 0


async def test_pdf_reader_metadata() -> None:
    block = PdfReaderBlock()
    result = await block.execute(PdfReaderInput(content=_b64("sample.pdf")))

    # metadata dict should exist (may or may not have keys depending on the PDF)
    assert isinstance(result.metadata, dict)


# ---------------------------------------------------------------------------
# PDF Renderer
# ---------------------------------------------------------------------------


async def test_pdf_renderer_page_1() -> None:
    block = PdfRendererBlock()
    result = await block.execute(
        PdfRendererInput(content=_b64("sample.pdf"), pages=[1])
    )

    assert result.page_count == 7
    assert len(result.images) == 1
    img = result.images[0]
    assert img["page_number"] == 1
    assert img["width"] > 0
    assert img["height"] > 0
    # Verify the image field is valid base64
    decoded = base64.b64decode(img["image"])
    assert len(decoded) > 100


async def test_pdf_renderer_dpi_72() -> None:
    block = PdfRendererBlock()
    result_72 = await block.execute(
        PdfRendererInput(content=_b64("sample.pdf"), pages=[1], dpi=72)
    )
    result_150 = await block.execute(
        PdfRendererInput(content=_b64("sample.pdf"), pages=[1], dpi=150)
    )

    # Higher DPI should produce larger images
    assert result_150.images[0]["width"] > result_72.images[0]["width"]
    assert result_150.images[0]["height"] > result_72.images[0]["height"]


async def test_pdf_renderer_png_format() -> None:
    block = PdfRendererBlock()
    result = await block.execute(
        PdfRendererInput(content=_b64("sample.pdf"), pages=[1], format="png")
    )

    decoded = base64.b64decode(result.images[0]["image"])
    # PNG files start with the PNG signature
    assert decoded[:4] == b"\x89PNG"


async def test_pdf_renderer_jpeg_format() -> None:
    block = PdfRendererBlock()
    result = await block.execute(
        PdfRendererInput(content=_b64("sample.pdf"), pages=[1], format="jpeg")
    )

    decoded = base64.b64decode(result.images[0]["image"])
    # JPEG files start with FF D8
    assert decoded[:2] == b"\xff\xd8"


# ---------------------------------------------------------------------------
# Excel Reader
# ---------------------------------------------------------------------------


async def test_excel_reader_known_data() -> None:
    block = ExcelReaderBlock()
    result = await block.execute(ExcelReaderInput(content=_b64("known_data.xlsx")))

    assert result.row_count == 3
    assert result.columns == ["Product", "Quantity", "Price"]
    assert result.records[0]["Product"] == "Widget A"
    assert result.records[0]["Quantity"] == 10
    assert result.records[0]["Price"] == 24.95


async def test_excel_reader_known_data_metadata_sheet() -> None:
    block = ExcelReaderBlock()
    result = await block.execute(
        ExcelReaderInput(content=_b64("known_data.xlsx"), sheet_name="Metadata")
    )

    assert result.row_count == 2
    assert result.columns == ["Key", "Value"]
    assert result.records[0]["Key"] == "Region"
    assert result.records[0]["Value"] == "Europe"


async def test_excel_reader_known_data_sheet_names() -> None:
    block = ExcelReaderBlock()
    result = await block.execute(ExcelReaderInput(content=_b64("known_data.xlsx")))

    assert "Sales" in result.sheet_names
    assert "Metadata" in result.sheet_names


async def test_excel_reader_sample_xlsx() -> None:
    block = ExcelReaderBlock()
    result = await block.execute(ExcelReaderInput(content=_b64("sample.xlsx")))

    # Should load without error and return some data
    assert isinstance(result.records, list)
    assert isinstance(result.columns, list)


# ---------------------------------------------------------------------------
# Excel Writer
# ---------------------------------------------------------------------------


async def test_excel_writer_valid_xlsx_output() -> None:
    block = ExcelWriterBlock()
    records: list[dict[str, Any]] = [
        {"Product": "Widget A", "Quantity": 10, "Price": 24.95},
        {"Product": "Widget B", "Quantity": 5, "Price": 49.95},
    ]
    result = await block.execute(ExcelWriterInput(records=records))

    assert result.row_count == 2
    assert result.column_count == 3

    raw = base64.b64decode(result.content)
    wb = openpyxl.load_workbook(io.BytesIO(raw))
    ws = wb.active
    assert ws is not None
    rows = list(ws.iter_rows(values_only=True))
    assert rows[0] == ("Product", "Quantity", "Price")
    assert rows[1] == ("Widget A", 10, 24.95)


async def test_excel_roundtrip() -> None:
    """Write records with ExcelWriter, read back with ExcelReader."""
    writer = ExcelWriterBlock()
    records: list[dict[str, Any]] = [
        {"Name": "Alice", "Score": 95},
        {"Name": "Bob", "Score": 87},
        {"Name": "Charlie", "Score": 73},
    ]
    write_result = await writer.execute(ExcelWriterInput(records=records))

    reader = ExcelReaderBlock()
    read_result = await reader.execute(
        ExcelReaderInput(content=write_result.content)
    )

    assert read_result.row_count == 3
    assert read_result.columns == ["Name", "Score"]
    assert read_result.records[0]["Name"] == "Alice"
    assert read_result.records[0]["Score"] == 95
    assert read_result.records[2]["Name"] == "Charlie"


# ---------------------------------------------------------------------------
# Word Reader
# ---------------------------------------------------------------------------


async def test_word_reader_known_data_text() -> None:
    block = WordReaderBlock()
    result = await block.execute(WordReaderInput(content=_b64("known_data.docx")))

    assert "Test Document" in result.text
    assert "first paragraph" in result.text


async def test_word_reader_known_data_paragraphs() -> None:
    block = WordReaderBlock()
    result = await block.execute(WordReaderInput(content=_b64("known_data.docx")))

    styles = [p["style"] for p in result.paragraphs]
    assert "Heading 1" in styles
    assert "Heading 2" in styles


async def test_word_reader_known_data_tables() -> None:
    block = WordReaderBlock()
    result = await block.execute(WordReaderInput(content=_b64("known_data.docx")))

    assert len(result.tables) >= 1
    table = result.tables[0]
    # Header row
    assert table[0] == ["Name", "Score"]
    # Data rows
    assert table[1] == ["Alice", "95"]
    assert table[2] == ["Bob", "87"]


async def test_word_reader_known_data_metadata() -> None:
    block = WordReaderBlock()
    result = await block.execute(WordReaderInput(content=_b64("known_data.docx")))

    assert isinstance(result.metadata, dict)


async def test_word_reader_sample_images_docx() -> None:
    block = WordReaderBlock()
    result = await block.execute(WordReaderInput(content=_b64("sample_images.docx")))

    # Should load without error
    assert isinstance(result.text, str)
    assert isinstance(result.paragraphs, list)


async def test_word_reader_sample_docx() -> None:
    block = WordReaderBlock()
    result = await block.execute(WordReaderInput(content=_b64("sample.docx")))

    assert isinstance(result.text, str)
    assert isinstance(result.paragraphs, list)


# ---------------------------------------------------------------------------
# Word Writer
# ---------------------------------------------------------------------------


async def test_word_writer_with_styles() -> None:
    block = WordWriterBlock()
    paragraphs = [
        {"text": "My Heading", "style": "Heading 1"},
        {"text": "Body text goes here.", "style": "Normal"},
        {"text": "A sub-heading", "style": "Heading 2"},
    ]
    result = await block.execute(WordWriterInput(paragraphs=paragraphs))

    assert result.paragraph_count == 3
    raw = base64.b64decode(result.content)
    doc = DocxDocument(io.BytesIO(raw))
    texts = [p.text for p in doc.paragraphs]
    assert "My Heading" in texts
    assert "Body text goes here." in texts


async def test_word_roundtrip() -> None:
    """Write paragraphs with WordWriter, read back with WordReader."""
    writer = WordWriterBlock()
    paragraphs = [
        {"text": "Roundtrip heading", "style": "Heading 1"},
        {"text": "Roundtrip body content."},
    ]
    write_result = await writer.execute(
        WordWriterInput(paragraphs=paragraphs, title="Roundtrip Doc")
    )

    reader = WordReaderBlock()
    read_result = await reader.execute(
        WordReaderInput(content=write_result.content)
    )

    assert "Roundtrip Doc" in read_result.text
    assert "Roundtrip heading" in read_result.text
    assert "Roundtrip body content." in read_result.text


# ---------------------------------------------------------------------------
# PowerPoint Reader
# ---------------------------------------------------------------------------


async def test_powerpoint_reader_known_data_slide_count() -> None:
    block = PowerpointReaderBlock()
    result = await block.execute(
        PowerpointReaderInput(content=_b64("known_data.pptx"))
    )

    assert result.slide_count == 3


async def test_powerpoint_reader_known_data_title() -> None:
    block = PowerpointReaderBlock()
    result = await block.execute(
        PowerpointReaderInput(content=_b64("known_data.pptx"))
    )

    slide_1_texts = result.slides[0]["texts"]
    assert any("Test Presentation" in t for t in slide_1_texts)


async def test_powerpoint_reader_known_data_findings() -> None:
    block = PowerpointReaderBlock()
    result = await block.execute(
        PowerpointReaderInput(content=_b64("known_data.pptx"))
    )

    slide_2_texts = result.slides[1]["texts"]
    assert any("Key Findings" in t for t in slide_2_texts)


async def test_powerpoint_reader_sample_pptx() -> None:
    block = PowerpointReaderBlock()
    result = await block.execute(
        PowerpointReaderInput(content=_b64("sample.pptx"))
    )

    assert result.slide_count >= 1
    assert isinstance(result.slides, list)


# ---------------------------------------------------------------------------
# PowerPoint Writer
# ---------------------------------------------------------------------------


async def test_powerpoint_writer_valid_output() -> None:
    block = PowerpointWriterBlock()
    slides = [
        {"title": "Slide One", "content": "Content A"},
        {"title": "Slide Two", "content": "Content B"},
    ]
    result = await block.execute(PowerpointWriterInput(slides=slides))

    assert result.slide_count == 2
    raw = base64.b64decode(result.content)
    prs = Presentation(io.BytesIO(raw))
    assert len(prs.slides) == 2


async def test_powerpoint_roundtrip() -> None:
    """Write slides with PowerpointWriter, read back with PowerpointReader."""
    writer = PowerpointWriterBlock()
    slides = [
        {"title": "RT Slide 1", "content": "Content 1"},
        {"title": "RT Slide 2", "content": "Content 2"},
    ]
    write_result = await writer.execute(
        PowerpointWriterInput(slides=slides, title="RT Presentation")
    )

    reader = PowerpointReaderBlock()
    read_result = await reader.execute(
        PowerpointReaderInput(content=write_result.content)
    )

    # title slide + 2 content slides
    assert read_result.slide_count == 3
    assert "RT Presentation" in read_result.text
    assert "RT Slide 1" in read_result.text


# ---------------------------------------------------------------------------
# YAML Parser
# ---------------------------------------------------------------------------


async def test_yaml_parser_known_data_project() -> None:
    block = YamlParserBlock()
    content = (FIXTURES / "known_data.yaml").read_text()
    result = await block.execute(YamlParserInput(content=content))

    assert result.data["project"] == "plumber"
    assert result.document_count == 1


async def test_yaml_parser_known_data_nested() -> None:
    block = YamlParserBlock()
    content = (FIXTURES / "known_data.yaml").read_text()
    result = await block.execute(YamlParserInput(content=content))

    assert result.data["config"]["mongo_uri"] == "mongodb://localhost:27017"
    assert result.data["config"]["redis_url"] == "redis://localhost:6379/0"


async def test_yaml_parser_known_data_features() -> None:
    block = YamlParserBlock()
    content = (FIXTURES / "known_data.yaml").read_text()
    result = await block.execute(YamlParserInput(content=content))

    assert result.data["features"] == ["pipelines", "triggers", "blocks"]


async def test_yaml_parser_multi_document() -> None:
    block = YamlParserBlock()
    content = (FIXTURES / "known_multi.yaml").read_text()
    result = await block.execute(
        YamlParserInput(content=content, multi_document=True)
    )

    assert result.document_count == 2
    assert result.data[0]["name"] == "pipeline-1"
    assert result.data[0]["blocks"] == 3
    assert result.data[1]["name"] == "pipeline-2"
    assert result.data[1]["blocks"] == 5


# ---------------------------------------------------------------------------
# YAML Writer
# ---------------------------------------------------------------------------


async def test_yaml_writer_roundtrip() -> None:
    block = YamlWriterBlock()
    original = {"project": "plumber", "version": "0.1.0", "active": True}
    result = await block.execute(YamlWriterInput(data=original))

    parsed = yaml.safe_load(result.content)
    assert parsed == original


async def test_yaml_writer_flow_style() -> None:
    block = YamlWriterBlock()
    data = {"x": 1, "y": 2}
    result_block = await block.execute(YamlWriterInput(data=data))
    result_flow = await block.execute(
        YamlWriterInput(data=data, default_flow_style=True)
    )

    # Flow style should be more compact (single line)
    assert len(result_flow.content.strip().splitlines()) < len(
        result_block.content.strip().splitlines()
    )


async def test_yaml_writer_list() -> None:
    block = YamlWriterBlock()
    data = ["alpha", "beta", "gamma"]
    result = await block.execute(YamlWriterInput(data=data))

    parsed = yaml.safe_load(result.content)
    assert parsed == data


# ---------------------------------------------------------------------------
# Parquet Reader
# ---------------------------------------------------------------------------


async def test_parquet_reader_known_data() -> None:
    block = ParquetReaderBlock()
    result = await block.execute(
        ParquetReaderInput(content=_b64("known_data.parquet"))
    )

    assert result.row_count == 5
    assert "id" in result.columns
    assert "name" in result.columns
    assert "value" in result.columns
    assert "active" in result.columns
    assert result.records[0]["name"] == "Alpha"
    assert result.records[4]["name"] == "Epsilon"


async def test_parquet_reader_select_columns() -> None:
    block = ParquetReaderBlock()
    result = await block.execute(
        ParquetReaderInput(
            content=_b64("known_data.parquet"), columns=["name", "value"]
        )
    )

    assert result.columns == ["name", "value"]
    assert "id" not in result.records[0]
    assert "active" not in result.records[0]
    assert result.records[0]["name"] == "Alpha"


async def test_parquet_reader_column_schema() -> None:
    block = ParquetReaderBlock()
    result = await block.execute(
        ParquetReaderInput(content=_b64("known_data.parquet"))
    )

    schema_by_name = {s["name"]: s["type"] for s in result.column_schema}
    assert "int" in schema_by_name["id"]
    assert "string" in schema_by_name["name"] or "utf8" in schema_by_name["name"]
    assert "double" in schema_by_name["value"] or "float" in schema_by_name["value"]
    assert "bool" in schema_by_name["active"]


# ---------------------------------------------------------------------------
# Parquet Writer
# ---------------------------------------------------------------------------


async def test_parquet_writer_valid_output() -> None:
    block = ParquetWriterBlock()
    records: list[dict[str, Any]] = [
        {"id": 1, "name": "Alpha", "value": 10.5},
        {"id": 2, "name": "Beta", "value": 20.3},
    ]
    result = await block.execute(ParquetWriterInput(records=records))

    assert result.row_count == 2
    assert result.column_count == 3

    raw = base64.b64decode(result.content)
    table = pq.read_table(io.BytesIO(raw))
    assert table.num_rows == 2


async def test_parquet_roundtrip() -> None:
    """Write records with ParquetWriter, read back with ParquetReader."""
    writer = ParquetWriterBlock()
    records: list[dict[str, Any]] = [
        {"id": 1, "name": "Alpha", "active": True},
        {"id": 2, "name": "Beta", "active": False},
        {"id": 3, "name": "Gamma", "active": True},
    ]
    write_result = await writer.execute(ParquetWriterInput(records=records))

    reader = ParquetReaderBlock()
    read_result = await reader.execute(
        ParquetReaderInput(content=write_result.content)
    )

    assert read_result.row_count == 3
    assert read_result.records[0]["name"] == "Alpha"
    assert read_result.records[1]["active"] is False
    assert read_result.records[2]["id"] == 3


# ---------------------------------------------------------------------------
# Markdown Renderer
# ---------------------------------------------------------------------------


async def test_markdown_renderer_known_data_headings() -> None:
    block = MarkdownRendererBlock()
    content = (FIXTURES / "known_data.md").read_text()
    result = await block.execute(MarkdownRendererInput(content=content))

    assert "<h1>" in result.html or "<h1" in result.html
    assert "Test Document" in result.html


async def test_markdown_renderer_known_data_table() -> None:
    block = MarkdownRendererBlock()
    content = (FIXTURES / "known_data.md").read_text()
    result = await block.execute(MarkdownRendererInput(content=content))

    assert "<table>" in result.html
    assert "Alpha" in result.html
    assert "Beta" in result.html


async def test_markdown_renderer_known_data_code() -> None:
    block = MarkdownRendererBlock()
    content = (FIXTURES / "known_data.md").read_text()
    result = await block.execute(MarkdownRendererInput(content=content))

    assert "<code" in result.html
    assert "def hello():" in result.text


async def test_markdown_renderer_text_strips_tags() -> None:
    block = MarkdownRendererBlock()
    content = (FIXTURES / "known_data.md").read_text()
    result = await block.execute(MarkdownRendererInput(content=content))

    assert "<h1>" not in result.text
    assert "<table>" not in result.text
    assert "<strong>" not in result.text
    # But the text content should still be present
    assert "Test Document" in result.text
    assert "Final paragraph." in result.text


async def test_markdown_renderer_bold_italic() -> None:
    block = MarkdownRendererBlock()
    content = (FIXTURES / "known_data.md").read_text()
    result = await block.execute(MarkdownRendererInput(content=content))

    assert "<strong>" in result.html
    assert "<em>" in result.html
    assert "bold" in result.text
    assert "italic" in result.text
