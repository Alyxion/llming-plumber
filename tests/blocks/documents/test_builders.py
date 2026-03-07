"""Tests for document builder blocks (excel, pdf, word, powerpoint)."""

from __future__ import annotations

import base64
import io
import json

from llming_plumber.blocks.documents.excel_builder import (
    CellStyle,
    ColumnDef,
    ExcelBuilderBlock,
    ExcelBuilderInput,
    MergedRange,
    SheetDef,
)
from llming_plumber.blocks.documents.pdf_builder import (
    FontSpec,
    LineElement,
    PageDef,
    PdfBuilderBlock,
    PdfBuilderInput,
    RectElement,
    TableDef,
    TextParagraph,
    TextRect,
)
from llming_plumber.blocks.documents.powerpoint_builder import (
    PowerpointBuilderBlock,
    PowerpointBuilderInput,
    ShapeDef,
    SlideTableDef,
    TextBoxDef,
    TextRun,
    TextStyle,
)
from llming_plumber.blocks.documents.powerpoint_builder import (
    SlideDef as PptxSlideDef,
)
from llming_plumber.blocks.documents.word_builder import (
    SectionDef,
    WordBuilderBlock,
    WordBuilderInput,
)

# ── Excel Builder ─────────────────────────────────────────────────────


async def test_excel_single_sheet() -> None:
    block = ExcelBuilderBlock()
    result = await block.execute(ExcelBuilderInput(
        sheets=[SheetDef(
            name="Data",
            rows=[
                {"name": "Alice", "score": 95},
                {"name": "Bob", "score": 87},
            ],
        )],
    ))
    assert result.sheet_count == 1
    assert result.total_rows == 2

    import openpyxl

    wb = openpyxl.load_workbook(
        io.BytesIO(base64.b64decode(result.content))
    )
    assert wb.sheetnames == ["Data"]
    ws = wb["Data"]
    assert ws.cell(1, 1).value == "name"
    assert ws.cell(2, 1).value == "Alice"
    assert ws.cell(3, 2).value == 87


async def test_excel_multi_sheet() -> None:
    block = ExcelBuilderBlock()
    result = await block.execute(ExcelBuilderInput(
        sheets=[
            SheetDef(
                name="Sales",
                columns=[
                    ColumnDef(key="product", header="Product", width=20),
                    ColumnDef(key="qty", header="Quantity"),
                ],
                rows=[{"product": "Widget", "qty": 10}],
            ),
            SheetDef(
                name="Summary",
                rows=[{"metric": "Total", "value": 42}],
            ),
        ],
    ))
    assert result.sheet_count == 2

    import openpyxl

    wb = openpyxl.load_workbook(
        io.BytesIO(base64.b64decode(result.content))
    )
    assert wb.sheetnames == ["Sales", "Summary"]
    assert wb["Sales"].cell(1, 1).value == "Product"
    assert wb["Summary"].cell(2, 1).value == "Total"


async def test_excel_styled_columns() -> None:
    block = ExcelBuilderBlock()
    result = await block.execute(ExcelBuilderInput(
        sheets=[SheetDef(
            name="Styled",
            columns=[
                ColumnDef(
                    key="amount",
                    header="Amount",
                    style=CellStyle(
                        number_format="#,##0.00", bold=True
                    ),
                    header_style=CellStyle(
                        bold=True, bg_color="FFFF00"
                    ),
                ),
            ],
            rows=[{"amount": 1234.5}],
        )],
    ))
    import openpyxl

    wb = openpyxl.load_workbook(
        io.BytesIO(base64.b64decode(result.content))
    )
    ws = wb["Styled"]
    assert ws.cell(1, 1).font.bold is True
    assert ws.cell(2, 1).number_format == "#,##0.00"


async def test_excel_freeze_and_filter() -> None:
    block = ExcelBuilderBlock()
    result = await block.execute(ExcelBuilderInput(
        sheets=[SheetDef(
            name="Frozen",
            rows=[{"a": 1}, {"a": 2}],
            freeze_panes="A2",
            auto_filter=True,
        )],
    ))
    import openpyxl

    wb = openpyxl.load_workbook(
        io.BytesIO(base64.b64decode(result.content))
    )
    ws = wb["Frozen"]
    assert ws.freeze_panes == "A2"
    assert ws.auto_filter.ref is not None


async def test_excel_merged_cells() -> None:
    block = ExcelBuilderBlock()
    result = await block.execute(ExcelBuilderInput(
        sheets=[SheetDef(
            name="Merged",
            columns=[
                ColumnDef(key="a"),
                ColumnDef(key="b"),
                ColumnDef(key="c"),
            ],
            rows=[{"a": 1, "b": 2, "c": 3}],
            merged_cells=[
                MergedRange(range="A1:C1", value="Title Row")
            ],
        )],
    ))
    import openpyxl

    wb = openpyxl.load_workbook(
        io.BytesIO(base64.b64decode(result.content))
    )
    ws = wb["Merged"]
    assert ws["A1"].value == "Title Row"


async def test_excel_from_json() -> None:
    block = ExcelBuilderBlock()
    sheets_json = json.dumps([{
        "name": "FromJSON",
        "rows": [{"x": 1}, {"x": 2}],
    }])
    result = await block.execute(ExcelBuilderInput(
        sheets=[SheetDef(name="unused")],
        json_data=sheets_json,
    ))
    assert result.sheet_count == 1
    assert result.total_rows == 2


async def test_excel_empty_sheet() -> None:
    block = ExcelBuilderBlock()
    result = await block.execute(ExcelBuilderInput(
        sheets=[SheetDef(name="Empty", rows=[])],
    ))
    assert result.total_rows == 0
    assert result.sheet_count == 1


# ── PDF Builder ───────────────────────────────────────────────────────


async def test_pdf_text_mode() -> None:
    block = PdfBuilderBlock()
    result = await block.execute(PdfBuilderInput(
        mode="text",
        pages=[PageDef(
            paragraphs=[
                TextParagraph(text="Main Title", style="heading1"),
                TextParagraph(
                    text="Some body text for the document."
                ),
                TextParagraph(text="A caption.", style="caption"),
            ],
        )],
        metadata={"title": "Test Doc", "author": "Plumber"},
    ))
    assert result.page_count == 1
    assert result.mode == "text"
    raw = base64.b64decode(result.content)
    assert raw[:5] == b"%PDF-"


async def test_pdf_geometric_mode() -> None:
    block = PdfBuilderBlock()
    result = await block.execute(PdfBuilderInput(
        mode="geometric",
        pages=[PageDef(
            text_rects=[
                TextRect(
                    x=72, y=750, text="Positioned text",
                    font=FontSpec(
                        size=18, bold=True, color="0000FF"
                    ),
                ),
                TextRect(
                    x=72, y=700, width=200,
                    text="Wrapped text in a box",
                    alignment="center",
                ),
            ],
            lines=[
                LineElement(
                    x1=72, y1=690, x2=500, y2=690,
                    color="FF0000",
                ),
            ],
            rects=[RectElement(
                x=72, y=600, width=200, height=80,
                fill_color="FFFFCC", stroke_color="000000",
            )],
        )],
    ))
    assert result.mode == "geometric"
    raw = base64.b64decode(result.content)
    assert raw[:5] == b"%PDF-"


async def test_pdf_mixed_mode() -> None:
    block = PdfBuilderBlock()
    result = await block.execute(PdfBuilderInput(
        mode="mixed",
        pages=[PageDef(
            paragraphs=[
                TextParagraph(text="Heading", style="heading1"),
            ],
            text_rects=[
                TextRect(x=400, y=750, text="Logo area"),
            ],
            lines=[
                LineElement(x1=72, y1=700, x2=523, y2=700),
            ],
        )],
    ))
    assert result.mode == "mixed"
    assert result.page_count == 1


async def test_pdf_multi_page() -> None:
    block = PdfBuilderBlock()
    result = await block.execute(PdfBuilderInput(
        mode="text",
        pages=[
            PageDef(
                paragraphs=[TextParagraph(text="Page 1")],
            ),
            PageDef(
                paragraphs=[TextParagraph(text="Page 2")],
            ),
            PageDef(
                width=842, height=595,
                paragraphs=[
                    TextParagraph(text="Page 3 landscape"),
                ],
            ),
        ],
    ))
    assert result.page_count == 3


async def test_pdf_table() -> None:
    block = PdfBuilderBlock()
    result = await block.execute(PdfBuilderInput(
        mode="geometric",
        pages=[PageDef(
            tables=[TableDef(
                x=72, y=700,
                headers=["Name", "Score"],
                rows=[["Alice", "95"], ["Bob", "87"]],
                col_widths=[200, 100],
                header_bg="CCCCCC",
            )],
        )],
    ))
    assert result.page_count == 1
    raw = base64.b64decode(result.content)
    assert raw[:5] == b"%PDF-"


async def test_pdf_from_json() -> None:
    block = PdfBuilderBlock()
    pages_json = json.dumps([{
        "paragraphs": [
            {"text": "From JSON", "style": "heading1"},
        ],
    }])
    result = await block.execute(PdfBuilderInput(
        mode="text",
        pages=[PageDef()],
        json_data=pages_json,
    ))
    assert result.page_count == 1


async def test_pdf_roundtrip_with_reader() -> None:
    """Build a PDF and read it back with PdfReaderBlock."""
    from llming_plumber.blocks.documents.pdf_reader import (
        PdfReaderBlock,
        PdfReaderInput,
    )

    builder = PdfBuilderBlock()
    built = await builder.execute(PdfBuilderInput(
        mode="text",
        pages=[PageDef(
            paragraphs=[
                TextParagraph(
                    text="Roundtrip Test", style="heading1"
                ),
                TextParagraph(
                    text="This text should survive the roundtrip."
                ),
            ],
        )],
    ))

    reader = PdfReaderBlock()
    read = await reader.execute(
        PdfReaderInput(content=built.content)
    )
    assert "Roundtrip Test" in read.text
    assert "roundtrip" in read.text.lower()


# ── Word Builder ──────────────────────────────────────────────────────


async def test_word_paragraphs_and_styles() -> None:
    block = WordBuilderBlock()
    result = await block.execute(WordBuilderInput(
        sections=[SectionDef(
            elements=[
                {
                    "type": "paragraph",
                    "text": "Title",
                    "style": "Heading 1",
                },
                {
                    "type": "paragraph",
                    "text": "Body text.",
                    "alignment": "justify",
                },
                {
                    "type": "paragraph",
                    "runs": [
                        {"text": "Bold ", "bold": True},
                        {
                            "text": "and italic",
                            "italic": True,
                            "font_color": "FF0000",
                        },
                    ],
                },
            ],
        )],
    ))
    assert result.section_count == 1
    assert result.element_count == 3

    import docx

    doc = docx.Document(
        io.BytesIO(base64.b64decode(result.content))
    )
    assert doc.paragraphs[0].text == "Title"


async def test_word_table() -> None:
    block = WordBuilderBlock()
    result = await block.execute(WordBuilderInput(
        sections=[SectionDef(
            elements=[{
                "type": "table",
                "rows": [
                    {
                        "cells": [
                            {
                                "text": "Name",
                                "bold": True,
                                "bg_color": "CCCCCC",
                            },
                            {
                                "text": "Score",
                                "bold": True,
                                "bg_color": "CCCCCC",
                            },
                        ],
                        "is_header": True,
                    },
                    {
                        "cells": [
                            {"text": "Alice"},
                            {"text": "95"},
                        ],
                    },
                    {
                        "cells": [
                            {"text": "Bob"},
                            {"text": "87"},
                        ],
                    },
                ],
            }],
        )],
    ))
    import docx

    doc = docx.Document(
        io.BytesIO(base64.b64decode(result.content))
    )
    assert len(doc.tables) == 1
    assert doc.tables[0].cell(1, 0).text == "Alice"


async def test_word_multi_section() -> None:
    block = WordBuilderBlock()
    result = await block.execute(WordBuilderInput(
        sections=[
            SectionDef(
                header_text="Section 1 Header",
                elements=[
                    {
                        "type": "paragraph",
                        "text": "First section.",
                    },
                ],
            ),
            SectionDef(
                footer_text="Page Footer",
                elements=[
                    {
                        "type": "paragraph",
                        "text": "Second section.",
                    },
                    {"type": "page_break"},
                    {
                        "type": "paragraph",
                        "text": "After break.",
                    },
                ],
            ),
        ],
    ))
    assert result.section_count == 2
    assert result.element_count == 4


async def test_word_from_json() -> None:
    block = WordBuilderBlock()
    sections_json = json.dumps([{
        "elements": [
            {"type": "paragraph", "text": "From JSON"},
        ],
    }])
    result = await block.execute(WordBuilderInput(
        sections=[SectionDef(elements=[])],
        json_data=sections_json,
    ))
    assert result.element_count == 1


async def test_word_roundtrip() -> None:
    """Build a docx and read it back."""
    from llming_plumber.blocks.documents.word_reader import (
        WordReaderBlock,
        WordReaderInput,
    )

    builder = WordBuilderBlock()
    built = await builder.execute(WordBuilderInput(
        sections=[SectionDef(
            elements=[
                {
                    "type": "paragraph",
                    "text": "Roundtrip Heading",
                    "style": "Heading 1",
                },
                {
                    "type": "paragraph",
                    "text": "Roundtrip body paragraph.",
                },
            ],
        )],
    ))

    reader = WordReaderBlock()
    read = await reader.execute(
        WordReaderInput(content=built.content)
    )
    assert "Roundtrip Heading" in read.text
    assert "Roundtrip body" in read.text


async def test_word_inline_formatting_roundtrip() -> None:
    block = WordBuilderBlock()
    result = await block.execute(WordBuilderInput(
        sections=[SectionDef(
            elements=[{
                "type": "paragraph",
                "runs": [
                    {"text": "Normal "},
                    {"text": "bold", "bold": True},
                    {"text": " and "},
                    {"text": "underlined", "underline": True},
                ],
            }],
        )],
    ))
    import docx

    doc = docx.Document(
        io.BytesIO(base64.b64decode(result.content))
    )
    runs = doc.paragraphs[0].runs
    assert any(r.bold for r in runs)
    assert any(r.underline for r in runs)


# ── PowerPoint Builder ────────────────────────────────────────────────


async def test_pptx_title_slide() -> None:
    block = PowerpointBuilderBlock()
    result = await block.execute(PowerpointBuilderInput(
        slides=[PptxSlideDef(
            layout="title",
            title="My Presentation",
            subtitle="By Plumber",
        )],
    ))
    assert result.slide_count == 1

    from pptx import Presentation

    prs = Presentation(
        io.BytesIO(base64.b64decode(result.content))
    )
    slide = prs.slides[0]
    assert slide.shapes.title.text == "My Presentation"


async def test_pptx_content_with_bullets() -> None:
    block = PowerpointBuilderBlock()
    result = await block.execute(PowerpointBuilderInput(
        slides=[PptxSlideDef(
            layout="title_and_content",
            title="Key Points",
            bullet_points=[
                "Point one", "Point two", "Point three",
            ],
        )],
    ))
    from pptx import Presentation

    prs = Presentation(
        io.BytesIO(base64.b64decode(result.content))
    )
    body = prs.slides[0].placeholders[1]
    texts = [p.text for p in body.text_frame.paragraphs]
    assert "Point one" in texts


async def test_pptx_text_box() -> None:
    block = PowerpointBuilderBlock()
    result = await block.execute(PowerpointBuilderInput(
        slides=[PptxSlideDef(
            layout="blank",
            text_boxes=[TextBoxDef(
                left=1, top=1, width=5, height=2,
                runs=[
                    TextRun(
                        text="Hello ",
                        style=TextStyle(bold=True),
                    ),
                    TextRun(
                        text="World",
                        style=TextStyle(font_color="FF0000"),
                    ),
                ],
            )],
        )],
    ))
    assert result.element_count >= 1


async def test_pptx_shapes() -> None:
    block = PowerpointBuilderBlock()
    result = await block.execute(PowerpointBuilderInput(
        slides=[PptxSlideDef(
            layout="blank",
            shapes=[
                ShapeDef(
                    shape_type="rectangle",
                    left=1, top=1, width=3, height=2,
                    fill_color="0066CC",
                    text="Box",
                    text_style=TextStyle(
                        bold=True, font_color="FFFFFF"
                    ),
                ),
                ShapeDef(
                    shape_type="oval",
                    left=5, top=1, width=2, height=2,
                    fill_color="CC0000",
                ),
            ],
        )],
    ))
    assert result.element_count >= 2


async def test_pptx_table() -> None:
    block = PowerpointBuilderBlock()
    result = await block.execute(PowerpointBuilderInput(
        slides=[PptxSlideDef(
            layout="blank",
            tables=[SlideTableDef(
                left=1, top=2, width=8, height=3,
                headers=["Name", "Score", "Grade"],
                rows=[
                    ["Alice", "95", "A"],
                    ["Bob", "87", "B"],
                ],
                header_bg="003366",
            )],
        )],
    ))
    from pptx import Presentation

    prs = Presentation(
        io.BytesIO(base64.b64decode(result.content))
    )
    shapes = list(prs.slides[0].shapes)
    table_shapes = [s for s in shapes if s.has_table]
    assert len(table_shapes) == 1
    tbl = table_shapes[0].table
    assert tbl.cell(0, 0).text == "Name"
    assert tbl.cell(1, 0).text == "Alice"


async def test_pptx_speaker_notes() -> None:
    block = PowerpointBuilderBlock()
    result = await block.execute(PowerpointBuilderInput(
        slides=[PptxSlideDef(
            layout="title_only",
            title="With Notes",
            notes="Remember to mention the key findings.",
        )],
    ))
    from pptx import Presentation

    prs = Presentation(
        io.BytesIO(base64.b64decode(result.content))
    )
    notes = prs.slides[0].notes_slide.notes_text_frame.text
    assert "key findings" in notes


async def test_pptx_multi_slide() -> None:
    block = PowerpointBuilderBlock()
    result = await block.execute(PowerpointBuilderInput(
        slides=[
            PptxSlideDef(layout="title", title="Slide 1"),
            PptxSlideDef(
                layout="title_and_content",
                title="Slide 2",
                bullet_points=["A", "B"],
            ),
            PptxSlideDef(layout="blank"),
        ],
    ))
    assert result.slide_count == 3


async def test_pptx_bg_color() -> None:
    block = PowerpointBuilderBlock()
    result = await block.execute(PowerpointBuilderInput(
        slides=[PptxSlideDef(
            layout="blank",
            bg_color="1A1A2E",
            text_boxes=[TextBoxDef(
                left=2, top=3, width=6, height=1,
                runs=[TextRun(
                    text="White on dark",
                    style=TextStyle(
                        font_color="FFFFFF", font_size=24
                    ),
                )],
            )],
        )],
    ))
    assert result.slide_count == 1


async def test_pptx_from_json() -> None:
    block = PowerpointBuilderBlock()
    slides_json = json.dumps([{
        "layout": "title",
        "title": "From JSON",
    }])
    result = await block.execute(PowerpointBuilderInput(
        slides=[PptxSlideDef(layout="blank")],
        json_data=slides_json,
    ))
    assert result.slide_count == 1


async def test_pptx_roundtrip() -> None:
    """Build a pptx and read it back."""
    from llming_plumber.blocks.documents.powerpoint_reader import (
        PowerpointReaderBlock,
        PowerpointReaderInput,
    )

    builder = PowerpointBuilderBlock()
    built = await builder.execute(PowerpointBuilderInput(
        slides=[
            PptxSlideDef(
                layout="title", title="Roundtrip Title"
            ),
            PptxSlideDef(
                layout="title_and_content",
                title="Content Slide",
                bullet_points=["Bullet A", "Bullet B"],
            ),
        ],
    ))

    reader = PowerpointReaderBlock()
    read = await reader.execute(
        PowerpointReaderInput(content=built.content)
    )
    assert read.slide_count == 2
    assert "Roundtrip Title" in read.text


# ── JSON serialization roundtrip ──────────────────────────────────────


def test_excel_sheet_def_json_roundtrip() -> None:
    """SheetDef serializes to JSON and back."""
    sheet = SheetDef(
        name="Test",
        columns=[
            ColumnDef(
                key="a", header="A", width=15,
                style=CellStyle(bold=True),
            ),
        ],
        rows=[{"a": 1}],
        freeze_panes="A2",
        auto_filter=True,
    )
    data = json.loads(sheet.model_dump_json())
    restored = SheetDef.model_validate(data)
    assert restored.name == "Test"
    assert restored.columns[0].style.bold is True
    assert restored.freeze_panes == "A2"


def test_pdf_page_def_json_roundtrip() -> None:
    """PageDef serializes to JSON and back."""
    page = PageDef(
        paragraphs=[
            TextParagraph(text="Hello", style="heading1"),
        ],
        text_rects=[
            TextRect(
                x=72, y=700, text="Pos",
                font=FontSpec(size=18),
            ),
        ],
        lines=[LineElement(x1=0, y1=0, x2=100, y2=100)],
        rects=[
            RectElement(
                x=10, y=10, width=50, height=50,
                fill_color="FF0000",
            ),
        ],
    )
    data = json.loads(page.model_dump_json())
    restored = PageDef.model_validate(data)
    assert restored.paragraphs[0].text == "Hello"
    assert restored.text_rects[0].font.size == 18
    assert restored.rects[0].fill_color == "FF0000"


def test_word_section_def_json_roundtrip() -> None:
    """SectionDef serializes to JSON and back."""
    section = SectionDef(
        elements=[
            {
                "type": "paragraph",
                "text": "Hi",
                "style": "Heading 1",
            },
            {
                "type": "table",
                "rows": [
                    {
                        "cells": [
                            {"text": "A"},
                            {"text": "B"},
                        ],
                    },
                ],
            },
        ],
        header_text="Header",
    )
    data = json.loads(section.model_dump_json())
    restored = SectionDef.model_validate(data)
    assert restored.elements[0]["text"] == "Hi"
    assert restored.header_text == "Header"


def test_pptx_slide_def_json_roundtrip() -> None:
    """SlideDef serializes to JSON and back."""
    slide = PptxSlideDef(
        layout="title_and_content",
        title="Test",
        bullet_points=["A", "B"],
        shapes=[ShapeDef(
            shape_type="oval",
            left=1, top=1, width=2, height=2,
            fill_color="FF0000", text="Circle",
        )],
        notes="Speaker notes here",
    )
    data = json.loads(slide.model_dump_json())
    restored = PptxSlideDef.model_validate(data)
    assert restored.title == "Test"
    assert restored.shapes[0].shape_type == "oval"
    assert restored.notes == "Speaker notes here"
