"""Generate test fixture files with known content.

Run once: python tests/fixtures/generate_fixtures.py

Sources and licenses:
- PDF: pdfplumber repo (MIT) — downloaded separately
- DOCX: python-openxml/python-docx repo (MIT) — downloaded separately
- PPTX: scanny/python-pptx repo (MIT) — downloaded separately
- XLSX: fluidware/openpyxl repo (MIT) — downloaded separately
- XLSX (known): generated here with openpyxl (MIT)
- DOCX (known): generated here with python-docx (MIT)
- PPTX (known): generated here with python-pptx (MIT)
- Parquet: generated here with pyarrow (Apache 2.0)
- YAML: plain text, no license concern
- Markdown: plain text, no license concern
"""

from __future__ import annotations

import json
from pathlib import Path

HERE = Path(__file__).parent


def make_xlsx() -> None:
    """Create an xlsx with known data for deterministic testing."""
    from openpyxl import Workbook

    wb = Workbook()

    # Sheet 1: Sales data
    ws1 = wb.active
    ws1.title = "Sales"
    ws1.append(["Product", "Quantity", "Price"])
    ws1.append(["Widget A", 10, 24.95])
    ws1.append(["Widget B", 5, 49.95])
    ws1.append(["Widget C", 20, 9.99])

    # Sheet 2: Metadata
    ws2 = wb.create_sheet("Metadata")
    ws2.append(["Key", "Value"])
    ws2.append(["Region", "Europe"])
    ws2.append(["Quarter", "Q1-2026"])

    wb.save(HERE / "known_data.xlsx")
    print("  Created known_data.xlsx (2 sheets, 3+2 rows)")


def make_docx() -> None:
    """Create a docx with known paragraphs and a table."""
    from docx import Document

    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is the first paragraph of the test document.")
    doc.add_heading("Section Two", level=2)
    doc.add_paragraph("Second paragraph with some content for testing.")

    # Add a table
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Score"
    table.cell(1, 0).text = "Alice"
    table.cell(1, 1).text = "95"
    table.cell(2, 0).text = "Bob"
    table.cell(2, 1).text = "87"

    doc.add_paragraph("Final paragraph.")

    doc.save(HERE / "known_data.docx")
    print("  Created known_data.docx (headings, paragraphs, table)")


def make_pptx() -> None:
    """Create a pptx with known slides."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Slide 1: Title
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test Presentation"
    slide.placeholders[1].text = "Plumber Integration Test"

    # Slide 2: Content
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Key Findings"
    body = slide.placeholders[1]
    body.text = "Finding one: data pipelines work."
    body.text_frame.add_paragraph().text = (
        "Finding two: triggers fire correctly."
    )
    body.text_frame.add_paragraph().text = (
        "Finding three: all formats supported."
    )

    # Slide 3: just a title
    layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(layout)

    prs.save(HERE / "known_data.pptx")
    print("  Created known_data.pptx (3 slides)")


def make_parquet() -> None:
    """Create a parquet file with known data."""
    import pyarrow as pa
    import pyarrow.parquet as pq

    table = pa.table({
        "id": [1, 2, 3, 4, 5],
        "name": ["Alpha", "Beta", "Gamma", "Delta", "Epsilon"],
        "value": [10.5, 20.3, 30.1, 40.7, 50.9],
        "active": [True, False, True, True, False],
    })
    pq.write_table(table, HERE / "known_data.parquet")
    print("  Created known_data.parquet (5 rows, 4 columns)")


def make_yaml() -> None:
    """Create YAML test fixtures."""
    import yaml

    single = {
        "project": "plumber",
        "version": "0.1.0",
        "features": ["pipelines", "triggers", "blocks"],
        "config": {
            "mongo_uri": "mongodb://localhost:27017",
            "redis_url": "redis://localhost:6379/0",
        },
    }
    (HERE / "known_data.yaml").write_text(
        yaml.dump(single, default_flow_style=False, sort_keys=False),
    )

    # Multi-document YAML
    docs = [
        {"name": "pipeline-1", "blocks": 3},
        {"name": "pipeline-2", "blocks": 5},
    ]
    (HERE / "known_multi.yaml").write_text(
        yaml.dump_all(docs, default_flow_style=False, sort_keys=False),
    )
    print("  Created known_data.yaml + known_multi.yaml")


def make_markdown() -> None:
    """Create a Markdown test fixture."""
    md = """\
# Test Document

This is a paragraph with **bold** and *italic* text.

## Data Table

| Name  | Value |
|-------|-------|
| Alpha | 10    |
| Beta  | 20    |

## Code Example

```python
def hello():
    print("Hello, World!")
```

Final paragraph.
"""
    (HERE / "known_data.md").write_text(md)
    print("  Created known_data.md")


def make_json() -> None:
    """Create a JSON test fixture."""
    data = {
        "pipelines": [
            {"name": "News Digest", "blocks": 2, "active": True},
            {"name": "Weather Monitor", "blocks": 2, "active": True},
            {"name": "Sentiment Tracker", "blocks": 3, "active": False},
        ],
        "version": "0.1.0",
    }
    (HERE / "known_data.json").write_text(
        json.dumps(data, indent=2) + "\n",
    )
    print("  Created known_data.json")


def make_csv() -> None:
    """Create a CSV test fixture."""
    csv_content = (
        "id,name,value,active\n"
        "1,Alpha,10.5,true\n"
        "2,Beta,20.3,false\n"
        "3,Gamma,30.1,true\n"
    )
    (HERE / "known_data.csv").write_text(csv_content)
    print("  Created known_data.csv")


if __name__ == "__main__":
    print("Generating test fixtures...")
    make_xlsx()
    make_docx()
    make_pptx()
    make_parquet()
    make_yaml()
    make_markdown()
    make_json()
    make_csv()
    print("Done!")
